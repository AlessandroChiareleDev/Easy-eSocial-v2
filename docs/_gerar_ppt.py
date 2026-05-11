"""Gera apresentacao PowerPoint sobre arquitetura multi-tenant do Easy-eSocial V2."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === Paleta ===
DARK_BG      = RGBColor(0x0B, 0x1E, 0x3A)   # azul escuro fundo
ACCENT       = RGBColor(0x4A, 0xDE, 0x80)   # verde menta
ACCENT_DARK  = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_TEXT   = RGBColor(0xF1, 0xF5, 0xF9)
MUTED_TEXT   = RGBColor(0xCB, 0xD5, 0xE1)
BLUE_BOX     = RGBColor(0x1E, 0x3A, 0x8A)
RED_ALERT    = RGBColor(0xEF, 0x44, 0x44)
ORANGE       = RGBColor(0xF9, 0x73, 0x16)
CARD_BG      = RGBColor(0x13, 0x2A, 0x4E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def fill_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=LIGHT_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=CARD_BG, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    s.shadow.inherit = False
    return s


def header(slide, title, subtitle=None):
    # barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6),
             title, size=28, bold=True, color=LIGHT_TEXT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=MUTED_TEXT)
    # linha sob titulo
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5),
                                 Inches(1.3) if subtitle else Inches(1.0),
                                 Inches(1.2), Emu(20000))
    ln.line.fill.background()
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "Easy-eSocial V2 — Arquitetura Multi-Tenant", size=10, color=MUTED_TEXT)
    add_text(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
             f"{page}", size=10, color=ACCENT, align=PP_ALIGN.RIGHT, bold=True)


# ========================================================================
# SLIDE 1 — CAPA
# ========================================================================
s = prs.slides.add_slide(BLANK)
fill_bg(s)

# faixas decorativas
for i, (yo, h, c) in enumerate([
    (Inches(0), Inches(0.15), ACCENT),
    (Inches(7.35), Inches(0.15), ACCENT),
]):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, yo, SW, h)
    b.line.fill.background(); b.fill.solid(); b.fill.fore_color.rgb = c

# circulo decorativo
circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.8), Inches(2.5), Inches(2.5))
circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT_DARK
circ.line.fill.background()
circ2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(1.4), Inches(1.4), Inches(1.4))
circ2.fill.solid(); circ2.fill.fore_color.rgb = DARK_BG
circ2.line.fill.background()

add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
         "EASY-ESOCIAL V2", size=22, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.4),
         "Arquitetura Multi-Tenant", size=54, bold=True, color=LIGHT_TEXT)
add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.7),
         "Como nosso sistema atende várias empresas com\num backend, dois Supabase e schemas isolados",
         size=20, color=MUTED_TEXT)

# tags
tags = ["2 empresas hoje", "2 Supabase", "Schema-per-tenant", "Pronto para escalar"]
xx = Inches(0.7)
for t in tags:
    w = Inches(0.15 * len(t) + 0.6)
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, xx, Inches(5.6), w, Inches(0.45))
    tag.fill.solid(); tag.fill.fore_color.rgb = CARD_BG
    tag.line.color.rgb = ACCENT; tag.line.width = Pt(1.25)
    add_text(s, xx, Inches(5.6), w, Inches(0.45), t, size=12, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    xx += w + Inches(0.15)

add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "Apresentação técnica  •  Maio/2026", size=14, color=MUTED_TEXT)

# ========================================================================
# SLIDE 2 — Agenda
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Agenda", "O que vamos cobrir nessa apresentação")

agenda = [
    ("01", "Panorama atual",       "Quantas empresas, quantos bancos, status de hoje"),
    ("02", "Modelo multi-tenant",  "Schema-per-empresa vs. DB-per-empresa"),
    ("03", "Os dois Supabase",     "Sistema DB e Dados DB — o que cada um faz"),
    ("04", "Backend em ação",      "Como uma request resolve a empresa correta"),
    ("05", "Onboarding",           "Como uma empresa nova entra no sistema"),
    ("06", "Escalabilidade",       "E se chegarem 10 empresas amanhã?"),
    ("07", "Roadmap & gaps",       "O que falta automatizar"),
]

y0 = Inches(1.7)
for i, (n, t, d) in enumerate(agenda):
    yy = y0 + Inches(0.72) * i
    # numero
    nb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), yy, Inches(0.55), Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb = ACCENT
    nb.line.fill.background()
    add_text(s, Inches(0.8), yy, Inches(0.55), Inches(0.55), n,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.6), yy + Inches(0.02), Inches(4.5), Inches(0.4),
             t, size=18, bold=True, color=LIGHT_TEXT)
    add_text(s, Inches(1.6), yy + Inches(0.36), Inches(11), Inches(0.4),
             d, size=13, color=MUTED_TEXT)
footer(s, 2)

# ========================================================================
# SLIDE 3 — Panorama (números)
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Panorama atual", "Estado de produção em 11/05/2026")

cards = [
    ("2",  "EMPRESAS ATIVAS",        "APPA + SOLUCOES",      ACCENT),
    ("2",  "PROJETOS SUPABASE",      "Sistema + Dados",       BLUE_BOX),
    ("36", "TABELAS POR EMPRESA",    "Schema completo eSocial", ORANGE),
    ("1",  "BACKEND FASTAPI",        "Multi-tenant nativo",   ACCENT_DARK),
]
cw = Inches(2.85); cgap = Inches(0.2); cx = Inches(0.7); cy = Inches(1.8)
for i, (big, lab, sub, col) in enumerate(cards):
    x = cx + (cw + cgap) * i
    box = add_box(s, x, cy, cw, Inches(2.4), fill=CARD_BG)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cy, cw, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, x, cy + Inches(0.3), cw, Inches(1.0), big,
             size=72, bold=True, color=col, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, cy + Inches(1.5), cw, Inches(0.4), lab,
             size=12, bold=True, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x, cy + Inches(1.9), cw, Inches(0.4), sub,
             size=11, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# linha embaixo: empresas detalhe
add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
         "EMPRESAS EM PRODUÇÃO", size=14, bold=True, color=ACCENT)

emp = [
    ("APPA",     "05969071000110", "schema: appa",     "v1.1.0", "376k eventos"),
    ("SOLUCOES", "09445502000109", "schema: solucoes", "v1.1.0", "schema novo"),
]
for i, (nome, cnpj, sch, ver, vol) in enumerate(emp):
    y = Inches(5.0) + Inches(0.85) * i
    box = add_box(s, Inches(0.7), y, Inches(12), Inches(0.75), fill=CARD_BG)
    add_text(s, Inches(0.9), y, Inches(2.5), Inches(0.75), nome,
             size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.4), y, Inches(2.5), Inches(0.75), cnpj,
             size=13, color=LIGHT_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.0), y, Inches(2.5), Inches(0.75), sch,
             size=13, color=MUTED_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.5), y, Inches(1.8), Inches(0.75), ver,
             size=13, color=MUTED_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.3), y, Inches(2.3), Inches(0.75), vol,
             size=13, bold=True, color=ORANGE, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ========================================================================
# SLIDE 4 — Modelo: schema por empresa
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Modelo: schema por empresa", "Não é 1 banco por empresa — é 1 schema PostgreSQL por empresa")

# ilustração 2 colunas: rejeitado vs adotado
# coluna 1
add_text(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.45),
         "❌ NÃO USAMOS", size=16, bold=True, color=RED_ALERT)
add_text(s, Inches(0.7), Inches(1.95), Inches(5.8), Inches(0.45),
         "Um banco Postgres por empresa", size=18, bold=True, color=LIGHT_TEXT)

cons = [
    "Custo: 1 projeto Supabase / empresa (US$ 25+/mês cada)",
    "Operação: 1 pool de conexões para cada uma",
    "Migrations: rodar N vezes em N bancos",
    "Backup: cresce linearmente",
]
for i, t in enumerate(cons):
    y = Inches(2.55) + Inches(0.45) * i
    bullet = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.12),
                                 Inches(0.12), Inches(0.12))
    bullet.fill.solid(); bullet.fill.fore_color.rgb = RED_ALERT
    bullet.line.fill.background()
    add_text(s, Inches(1.1), y, Inches(5.4), Inches(0.4), t, size=13, color=MUTED_TEXT)

# coluna 2
add_text(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.45),
         "✅ ADOTAMOS", size=16, bold=True, color=ACCENT)
add_text(s, Inches(7.0), Inches(1.95), Inches(5.8), Inches(0.45),
         "Um schema PostgreSQL por empresa", size=18, bold=True, color=LIGHT_TEXT)
pros = [
    "1 só projeto Supabase para todos os tenants",
    "Pool único de conexões — eficiente",
    "Migrations: 1 SQL aplicado N vezes via runner",
    "Isolamento por search_path + checagem no backend",
    "Suporta dezenas de empresas sem mudar infra",
]
for i, t in enumerate(pros):
    y = Inches(2.55) + Inches(0.45) * i
    bullet = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), y + Inches(0.12),
                                 Inches(0.12), Inches(0.12))
    bullet.fill.solid(); bullet.fill.fore_color.rgb = ACCENT
    bullet.line.fill.background()
    add_text(s, Inches(7.4), y, Inches(5.4), Inches(0.4), t, size=13, color=MUTED_TEXT)

# nota destaque embaixo
note = add_box(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.2),
               fill=BLUE_BOX, line=ACCENT)
add_text(s, Inches(0.95), Inches(5.7), Inches(11.5), Inches(0.4),
         "💡 ANALOGIA", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.7),
         "Pense num condomínio: 1 prédio (banco), N apartamentos (schemas). Cada\n"
         "empresa tem chave que abre só o seu apartamento — search_path é a chave.",
         size=14, color=LIGHT_TEXT)
footer(s, 4)

# ========================================================================
# SLIDE 5 — Os dois Supabase (diagrama)
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Os dois Supabase", "Sistema DB (auth/routing) e Dados DB (schemas das empresas)")

# Caixa esquerda — Sistema
sys_box = add_box(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.0),
                  fill=CARD_BG, line=ACCENT)
add_text(s, Inches(0.9), Inches(1.75), Inches(5.4), Inches(0.4),
         "🔐  SISTEMA DB", size=18, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(2.15), Inches(5.4), Inches(0.4),
         "db.kjbgiwnlvqnrfdozjvhq.supabase.co", size=11, color=MUTED_TEXT,
         font="Consolas")
add_text(s, Inches(0.9), Inches(2.55), Inches(5.4), Inches(0.4),
         "Auth + Routing + Auditoria", size=13, color=LIGHT_TEXT)

tbls = [
    ("users",            "Credenciais bcrypt + flags"),
    ("empresas_routing", "CNPJ → schema_name (registry)"),
    ("user_empresas",    "N:N usuário ↔ empresa + papel"),
    ("audit_log",        "Log append-only de ações"),
    ("schema_meta",      "Versão do schema sistema"),
]
for i, (nm, dsc) in enumerate(tbls):
    y = Inches(3.15) + Inches(0.65) * i
    bx = add_box(s, Inches(0.95), y, Inches(5.3), Inches(0.55), fill=DARK_BG)
    add_text(s, Inches(1.05), y, Inches(2.2), Inches(0.55), nm,
             size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    add_text(s, Inches(3.05), y, Inches(3.15), Inches(0.55), dsc,
             size=11, color=MUTED_TEXT, anchor=MSO_ANCHOR.MIDDLE)

# Caixa direita — Dados
dat_box = add_box(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(5.0),
                  fill=CARD_BG, line=ORANGE)
add_text(s, Inches(7.05), Inches(1.75), Inches(5.6), Inches(0.4),
         "💾  DADOS DB", size=18, bold=True, color=ORANGE)
add_text(s, Inches(7.05), Inches(2.15), Inches(5.6), Inches(0.4),
         "aws-1-us-east-2.pooler.supabase.com", size=11, color=MUTED_TEXT, font="Consolas")
add_text(s, Inches(7.05), Inches(2.55), Inches(5.6), Inches(0.4),
         "Schemas das empresas (36 tabelas cada)", size=13, color=LIGHT_TEXT)

schemas = [
    ("appa",     "APPA — 376k eventos",  ACCENT),
    ("solucoes", "SOLUCOES — schema novo", ACCENT),
    ("legado",   "V1 raw dump (read-only)", MUTED_TEXT),
    ("public",   "Tabelas compartilhadas", MUTED_TEXT),
]
for i, (sch, dsc, col) in enumerate(schemas):
    y = Inches(3.15) + Inches(0.78) * i
    bx = add_box(s, Inches(7.1), y, Inches(5.5), Inches(0.68), fill=DARK_BG)
    add_text(s, Inches(7.25), y, Inches(2.0), Inches(0.68), sch,
             size=13, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    add_text(s, Inches(9.0), y, Inches(3.5), Inches(0.68), dsc,
             size=11, color=MUTED_TEXT, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 5)

# ========================================================================
# SLIDE 6 — Fluxo de uma request
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Como uma request resolve a empresa", "Da tela do usuário até o schema correto, em milissegundos")

steps = [
    ("1", "FRONTEND",      "Vue envia request:\nAuthorization: Bearer <JWT>\nX-Empresa-CNPJ: 05969071000110", ACCENT),
    ("2", "NGINX",         "Recebe HTTPS 443\nProxy → 127.0.0.1:8001", BLUE_BOX),
    ("3", "FASTAPI",       "Decodifica JWT\nValida user_empresas\nResolve CNPJ → schema", ORANGE),
    ("4", "POSTGRES",      "SET search_path TO\n  appa, public\nExecuta query", ACCENT_DARK),
    ("5", "RESPONSE",      "JSON volta isolado\nRESET search_path\nConn devolvida ao pool", ACCENT),
]
n = len(steps)
total_w = Inches(12.6)
gap = Inches(0.15)
bw = (total_w - gap * (n-1)) / n
y0 = Inches(2.2)

for i, (num, tit, dsc, col) in enumerate(steps):
    x = Inches(0.7) + (bw + gap) * i
    box = add_box(s, x, y0, bw, Inches(3.4), fill=CARD_BG)
    # número grande
    nb = s.shapes.add_shape(MSO_SHAPE.OVAL, x + bw/2 - Inches(0.4), y0 + Inches(0.2),
                             Inches(0.8), Inches(0.8))
    nb.fill.solid(); nb.fill.fore_color.rgb = col; nb.line.fill.background()
    add_text(s, x + bw/2 - Inches(0.4), y0 + Inches(0.2), Inches(0.8), Inches(0.8),
             num, size=24, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y0 + Inches(1.15), bw, Inches(0.45), tit,
             size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y0 + Inches(1.65), bw - Inches(0.2), Inches(1.7), dsc,
             size=11, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

    # seta entre caixas
    if i < n - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    x + bw + Inches(0.005), y0 + Inches(1.5),
                                    gap - Inches(0.01), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# código embaixo
code_box = add_box(s, Inches(0.7), Inches(5.85), Inches(12.6), Inches(1.05), fill=RGBColor(0x05,0x10,0x20))
add_text(s, Inches(0.9), Inches(5.92), Inches(12.2), Inches(0.4),
         "# backend/app/tenant.py", size=10, color=MUTED_TEXT, font="Consolas")
add_text(s, Inches(0.9), Inches(6.25), Inches(12.2), Inches(0.6),
         'with empresa_conn(cnpj) as conn:  # SET search_path TO "appa", public',
         size=14, bold=True, color=ACCENT, font="Consolas")
footer(s, 6)

# ========================================================================
# SLIDE 7 — Onboarding (passo a passo)
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Onboarding de empresa nova", "4 passos — ~30 segundos no terminal")

steps = [
    ("1️⃣", "CRIAR SCHEMA",
     "CREATE SCHEMA acme_corp;",
     "No Dados DB, via psql admin"),
    ("2️⃣", "RODAR MIGRATION",
     "python -m app.migrate apply\n  --target empresa --version 1.1.0\n  --schema acme_corp",
     "Runner cria as 36 tabelas no schema novo"),
    ("3️⃣", "REGISTRAR EMPRESA",
     "INSERT INTO sistema.empresas_routing\n  (cnpj, razao_social, schema_name,\n   schema_version, ativo)\nVALUES (...)",
     "Vincula CNPJ → schema no registry"),
    ("4️⃣", "VINCULAR USUÁRIO",
     "INSERT INTO sistema.user_empresas\n  (user_id, cnpj, papel)\nVALUES (...)",
     "Define quem pode operar a empresa"),
]
for i, (icon, tit, code, dsc) in enumerate(steps):
    row, col = divmod(i, 2)
    x = Inches(0.7) + Inches(6.3) * col
    y = Inches(1.65) + Inches(2.6) * row
    box = add_box(s, x, y, Inches(6.0), Inches(2.4), fill=CARD_BG, line=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(0.7), Inches(0.6),
             icon, size=24, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.9), y + Inches(0.15), Inches(5.0), Inches(0.5),
             tit, size=16, bold=True, color=ACCENT)
    # bloco de código
    cb = add_box(s, x + Inches(0.2), y + Inches(0.75), Inches(5.6), Inches(1.15),
                 fill=RGBColor(0x05,0x10,0x20))
    add_text(s, x + Inches(0.35), y + Inches(0.8), Inches(5.4), Inches(1.05),
             code, size=11, color=ACCENT, font="Consolas")
    add_text(s, x + Inches(0.2), y + Inches(1.95), Inches(5.6), Inches(0.4),
             dsc, size=11, color=MUTED_TEXT)
footer(s, 7)

# ========================================================================
# SLIDE 8 — Escalabilidade (10 empresas)
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "E se chegarem 10 empresas amanhã?", "A arquitetura escala — só roda o script em loop")

# Gráfico simples horizontal de capacidade
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.45),
         "CAPACIDADE DO MODELO ATUAL", size=14, bold=True, color=ACCENT)

bars = [
    ("Hoje",                   2,  ACCENT,     "Cabe folgado"),
    ("+ 10 novas (próx. mês)", 12, ACCENT_DARK,"Sem ajuste"),
    ("+ 30 (1 ano)",           32, ORANGE,     "Subir pool p/ (5, 50)"),
    ("+ 100 (futuro)",        102, RED_ALERT,  "Pensar em sharding/DB-per-cliente"),
]
maxv = 110
chart_x = Inches(0.7); chart_w = Inches(8.5); barh = Inches(0.55); ygap = Inches(0.85)
y0 = Inches(2.3)
for i, (lab, val, col, note) in enumerate(bars):
    y = y0 + ygap * i
    # trilho
    track = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, chart_x, y, chart_w, barh)
    track.fill.solid(); track.fill.fore_color.rgb = CARD_BG; track.line.fill.background()
    # barra
    bw = int(chart_w * (val / maxv))
    if bw < Inches(0.5): bw = Inches(0.5)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, chart_x, y, bw, barh)
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, chart_x + Inches(0.15), y, bw, barh, lab,
             size=12, bold=True, color=DARK_BG, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, chart_x + chart_w + Inches(0.15), y, Inches(0.8), barh, f"{val}",
             size=14, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, chart_x + chart_w + Inches(0.9), y, Inches(3.5), barh, note,
             size=11, color=MUTED_TEXT, anchor=MSO_ANCHOR.MIDDLE)

# painel de custo embaixo
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
         "CUSTO POR EMPRESA NOVA",
         size=14, bold=True, color=ACCENT)

costs = [
    ("⏱️ Tempo", "~30 segundos de comandos"),
    ("💾 Espaço", "~50 MB inicial (schema vazio)"),
    ("💰 Dinheiro", "R$ 0 adicional (mesmo banco)"),
    ("👷 Pessoas", "1 operador com acesso SSH"),
]
cw = Inches(2.95); cx = Inches(0.7); cy = Inches(6.45)
for i, (lab, val) in enumerate(costs):
    x = cx + (cw + Inches(0.1)) * i
    bx = add_box(s, x, cy, cw, Inches(0.55), fill=CARD_BG)
    add_text(s, x + Inches(0.15), cy, cw, Inches(0.55),
             f"{lab}  {val}", size=11, color=LIGHT_TEXT, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8)

# ========================================================================
# SLIDE 9 — Stack técnica
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Stack técnica", "O que está rodando em produção")

cats = [
    ("Frontend", ACCENT, [
        "Vue 3 + <script setup>",
        "Vite (build estático)",
        "Pinia (stores: auth, empresa)",
        "vue-router",
        "Servido pelo nginx em /opt/easy-esocial/frontend-dist/",
    ]),
    ("Backend", ORANGE, [
        "Python 3.12 + venv",
        "FastAPI 0.115.6 + uvicorn",
        "psycopg2-binary + ThreadedConnectionPool(2..20)",
        "JWT HS256 (8h) + bcrypt",
        "Fernet para senhas de certificado A1",
        "systemd: easy-esocial.service (port 8001)",
    ]),
    ("Banco de dados", ACCENT_DARK, [
        "PostgreSQL 15 (Supabase)",
        "Sistema DB: 1 schema (sistema), 5 tabelas",
        "Dados DB: schemas appa, solucoes (+legado, public)",
        "36 tabelas por empresa",
        "search_path por conexão (multi-tenant lógico)",
        "Migrations versionadas em backend/migrations/",
    ]),
    ("Infra", BLUE_BOX, [
        "VPS Hostinger (76.13.169.45)",
        "Ubuntu + nginx (TLS Let's Encrypt)",
        "Deploy: GitHub Action SSH → /opt/easy-esocial/repo",
        "deploy.sh + systemctl restart (~60s)",
        "Backup manual via pg_dump (a melhorar)",
    ]),
]
for i, (cat, col, items) in enumerate(cats):
    row, ci = divmod(i, 2)
    x = Inches(0.7) + Inches(6.3) * ci
    y = Inches(1.6) + Inches(2.7) * row
    box = add_box(s, x, y, Inches(6.0), Inches(2.55), fill=CARD_BG)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.15), Inches(2.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(5.5), Inches(0.45),
             cat, size=16, bold=True, color=col)
    for j, it in enumerate(items):
        yy = y + Inches(0.55) + Inches(0.32) * j
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), yy + Inches(0.1),
                                     Inches(0.1), Inches(0.1))
        bullet.fill.solid(); bullet.fill.fore_color.rgb = col; bullet.line.fill.background()
        add_text(s, x + Inches(0.55), yy, Inches(5.3), Inches(0.32),
                 it, size=11, color=LIGHT_TEXT)
footer(s, 9)

# ========================================================================
# SLIDE 10 — Roadmap / Gaps
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
header(s, "Próximos passos", "O que ainda falta — e o que protege bem")

# OK
add_text(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.45),
         "✅  O QUE JÁ ESTÁ SÓLIDO", size=16, bold=True, color=ACCENT)
ok = [
    "Routing dinâmico por CNPJ",
    "Migrations idempotentes versionadas",
    "JWT + bcrypt + Fernet em certificados",
    "RESET search_path em todo finally do pool",
    "Audit log de ações sensíveis",
    "Deploy automatizado via GitHub Action",
]
for i, t in enumerate(ok):
    y = Inches(2.2) + Inches(0.45) * i
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.12),
                            Inches(0.14), Inches(0.14))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background()
    add_text(s, Inches(1.1), y, Inches(5.3), Inches(0.4), t, size=13, color=LIGHT_TEXT)

# Pendente
add_text(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.45),
         "⚠️  PRÓXIMAS MELHORIAS", size=16, bold=True, color=ORANGE)
todo = [
    "UI no admin para cadastrar empresa (hoje é SQL)",
    "Endpoint POST /api/admin/users",
    "RLS (Row Level Security) no Dados DB",
    "Health check periódico no pool",
    "Retenção/rotação de audit_log",
    "Dashboard: requests/conexões por empresa",
    "Backup automatizado por schema",
]
for i, t in enumerate(todo):
    y = Inches(2.2) + Inches(0.45) * i
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), y + Inches(0.12),
                            Inches(0.14), Inches(0.14))
    b.fill.solid(); b.fill.fore_color.rgb = ORANGE; b.line.fill.background()
    add_text(s, Inches(7.4), y, Inches(5.3), Inches(0.4), t, size=13, color=LIGHT_TEXT)

# nota inferior
note = add_box(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.85),
               fill=BLUE_BOX, line=ACCENT)
add_text(s, Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.4),
         "📌 PRINCÍPIO", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.5),
         "Hoje a arquitetura aguenta dezenas de empresas. Migração para "
         "DB-per-empresa só faz sentido quando um cliente passar de 5GB ou exigir isolamento físico contratual.",
         size=12, color=LIGHT_TEXT)
footer(s, 10)

# ========================================================================
# SLIDE 11 — Encerramento
# ========================================================================
s = prs.slides.add_slide(BLANK); fill_bg(s)
# fundo decorativo
for i in range(6):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(10 + i*0.4), Inches(-0.5 + i*0.3),
                            Inches(3 - i*0.3), Inches(3 - i*0.3))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT_DARK if i%2 else BLUE_BOX
    c.line.fill.background()
    c.fill.fore_color.brightness = 0.0

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0), Inches(0.4), SH)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text(s, Inches(0.9), Inches(1.8), Inches(11), Inches(0.5),
         "EM RESUMO", size=18, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(2.4), Inches(11), Inches(1.3),
         "1 backend, 2 Supabase,\nN empresas — sem dor.",
         size=42, bold=True, color=LIGHT_TEXT)

bullets = [
    "📦  2 empresas hoje (APPA + SOLUCOES) — modelo já provado",
    "🏢  Schema-per-empresa: barato, escalável, simples de operar",
    "🚀  Onboarding em 4 passos, ~30s — só falta UI",
    "📈  Capacidade folgada para dezenas de novas empresas",
]
for i, t in enumerate(bullets):
    add_text(s, Inches(0.9), Inches(4.6) + Inches(0.42) * i, Inches(11), Inches(0.4),
             t, size=15, color=LIGHT_TEXT)

add_text(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.5),
         "Easy-eSocial V2  •  Maio/2026",
         size=12, color=MUTED_TEXT)

OUT = r"c:\Users\xandao\Documents\GitHub\Easy-eSocial-v2\docs\APRESENTACAO_ARQUITETURA_MULTI_TENANT.pptx"
prs.save(OUT)
print("OK ->", OUT)
